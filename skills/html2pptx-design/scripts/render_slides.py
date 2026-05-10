# -*- coding: utf-8 -*-
"""Render slides.html into TWO PowerPoint files:

  1. <output>.pptx          — pixel-perfect, image-based (1 PNG screenshot per slide)
  2. <output (-editable)>.pptx — text-editable, reconstructed from the DOM
                                 (real text boxes + rounded rectangles, positions
                                  and colours pulled from getComputedStyle)

Usage:
    python render_slides.py <slides_html_path> <output_pptx_path> [--total N] [--no-editable] [--only-editable]

Examples:
    python render_slides.py slides.html deck.pptx
        -> deck.pptx (image) + deck-editable.pptx (editable)
    python render_slides.py slides.html deck.pptx --no-editable
        -> deck.pptx only
"""

import asyncio
import argparse
import re
from pathlib import Path

from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# 1920px-wide slide mapped to 13.333in (= 144 px/in). So:
PX_TO_EMU = 6350  # 12192000 EMU / 1920 px
PX_TO_PT = 0.5  # 72 pt/in / 144 px/in
SLIDE_W_PX, SLIDE_H_PX = 1920, 1080
DEFAULT_BG = RGBColor(0x10, 0x18, 0x28)


# --------------------------------------------------------------------------- #
# DOM extraction (runs in the browser)
# --------------------------------------------------------------------------- #

_EXTRACT_JS = r"""
() => {
  const INLINE = new Set(['SPAN','B','I','EM','STRONG','CODE','A','SMALL','MARK','U','SUB','SUP','BR']);
  const isTransparent = (c) => !c || c === 'transparent' || c === 'rgba(0, 0, 0, 0)';
  const firstRgb = (s) => { if(!s) return null; const m = s.match(/rgba?\([^)]+\)/); return m ? m[0] : null; };
  const hasBlockChild = (el) => Array.from(el.children).some(ch => !INLINE.has(ch.tagName));

  const slides = [];
  let i = 1;
  while (document.querySelector('#s' + i)) {
    const root = document.querySelector('#s' + i);
    const rootRect = root.getBoundingClientRect();
    const rs = getComputedStyle(root);
    const slide = {
      w: rootRect.width, h: rootRect.height,
      bg: !isTransparent(rs.backgroundColor) ? rs.backgroundColor : (firstRgb(rs.backgroundImage) || null),
      boxes: [], texts: [],
    };
    const rel = (r) => ({ x: r.left - rootRect.left, y: r.top - rootRect.top, w: r.width, h: r.height });

    const collectRuns = (el) => {
      const runs = [];
      const walkInline = (node) => {
        node.childNodes.forEach(ch => {
          if (ch.nodeType === 3) {
            const t = ch.textContent;
            if (!t) return;
            const cs = getComputedStyle(node);
            runs.push({ text: t, size: parseFloat(cs.fontSize), color: cs.color,
                        bold: (parseInt(cs.fontWeight)||400) >= 600,
                        family: (cs.fontFamily||'').split(',')[0].replace(/['"]/g,'').trim(),
                        mono: /mono|consol|courier|jetbrains/i.test(cs.fontFamily||'') });
          } else if (ch.nodeType === 1) {
            if (ch.tagName === 'BR') { runs.push({ text: '\n', size: 0, color: '', bold:false, family:'', mono:false }); }
            else if (INLINE.has(ch.tagName)) { walkInline(ch); }
          }
        });
      };
      walkInline(el);
      // drop leading/trailing empty
      while (runs.length && !runs[0].text.replace(/\n/g,'').trim() && runs[0].text !== '\n') runs.shift();
      return runs;
    };

    const visit = (el) => {
      if (el !== root) {
        const st = getComputedStyle(el);
        if (st.display === 'none' || st.visibility === 'hidden' || parseFloat(st.opacity) === 0) return;
        const rect = el.getBoundingClientRect();
        const r = rel(rect);

        // -- box? (background colour, gradient, or visible border) --
        const bgCol = !isTransparent(st.backgroundColor) ? st.backgroundColor : null;
        const bgImg = (!bgCol && st.backgroundImage && st.backgroundImage !== 'none') ? firstRgb(st.backgroundImage) : null;
        const bw = Math.max(parseFloat(st.borderTopWidth)||0, parseFloat(st.borderLeftWidth)||0,
                            parseFloat(st.borderRightWidth)||0, parseFloat(st.borderBottomWidth)||0);
        const borderCol = (bw > 0.4 && !isTransparent(st.borderTopColor)) ? st.borderTopColor : null;
        const fill = bgCol || bgImg;
        if ((fill || borderCol) && r.w > 3 && r.h > 3 && r.w < 2000 && r.h < 1200) {
          const radius = Math.max(parseFloat(st.borderTopLeftRadius)||0, 0);
          const isOval = (st.borderRadius && /(50%|9999px|999px)/.test(st.borderRadius) && Math.abs(r.w - r.h) < Math.max(r.w,r.h)*0.25);
          slide.boxes.push({ x:r.x, y:r.y, w:r.w, h:r.h, fill, border:borderCol, borderW:bw||1,
                             radius: Math.min(radius, Math.min(r.w,r.h)/2), oval: !!isOval, area: r.w*r.h });
        }

        // -- text? (this element directly holds text and has no block-level children) --
        const directText = Array.from(el.childNodes).some(n => n.nodeType === 3 && n.textContent.trim());
        const onlyInline = !hasBlockChild(el);
        const innerTxt = (el.innerText || '').trim();
        if (innerTxt && (directText || (el.children.length > 0 && onlyInline)) && onlyInline) {
          const cs = getComputedStyle(el);
          let runs = collectRuns(el);
          if (!runs.length) runs = [{ text: innerTxt, size: parseFloat(cs.fontSize), color: cs.color,
                                      bold:(parseInt(cs.fontWeight)||400)>=600,
                                      family:(cs.fontFamily||'').split(',')[0].replace(/['"]/g,'').trim(),
                                      mono:/mono|consol|courier|jetbrains/i.test(cs.fontFamily||'') }];
          let lh = cs.lineHeight;
          let lineRatio = 1.15;
          const fpx = parseFloat(cs.fontSize) || 16;
          if (lh && lh !== 'normal') { const v = parseFloat(lh); lineRatio = lh.endsWith('px') ? v / fpx : v; }
          slide.texts.push({ x:r.x, y:r.y, w:r.w, h:r.h, align: cs.textAlign || 'left',
                             lineRatio: lineRatio, fontSize: fpx, runs });
          // recurse only into block children (none, since onlyInline) -> stop here
          return;
        }
      }
      Array.from(el.children).forEach(visit);
    };
    visit(root);
    slides.push(slide);
    i++;
  }
  return slides;
}
"""


def _rgb(css):
    """'rgb(r,g,b)' / 'rgba(r,g,b,a)' -> RGBColor (or None when fully transparent)."""
    if not css:
        return None
    m = re.match(
        r"rgba?\(\s*([\d.]+)\s*,\s*([\d.]+)\s*,\s*([\d.]+)\s*(?:,\s*([\d.]+)\s*)?\)",
        css.strip(),
    )
    if not m:
        return None
    r, g, b = int(float(m.group(1))), int(float(m.group(2))), int(float(m.group(3)))
    a = float(m.group(4)) if m.group(4) is not None else 1.0
    if a < 0.06:
        return None
    return RGBColor(max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))


def _align(css):
    return {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
        "start": PP_ALIGN.LEFT,
        "end": PP_ALIGN.RIGHT,
    }.get((css or "left").lower(), PP_ALIGN.LEFT)


def _set_ea(run, typeface):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", typeface)


# --------------------------------------------------------------------------- #
# Renderers
# --------------------------------------------------------------------------- #


async def render_screenshots(html_path: Path, out_dir: Path, total):
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        ctx = await browser.new_context(
            viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX}, device_scale_factor=1
        )
        page = await ctx.new_page()
        await page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
        await page.wait_for_timeout(1500)

        if total is None:
            total = 0
            while await page.query_selector(f"#s{total + 1}"):
                total += 1
            print(f"  auto-detected {total} slides")

        for i in range(1, total + 1):
            el = await page.query_selector(f"#s{i}")
            if not el:
                print(f"!! missing #s{i}")
                continue
            await el.screenshot(path=str(out_dir / f"slide-{i:02d}.png"))
            print(f"  rendered slide-{i:02d}.png")

        dom = await page.evaluate(_EXTRACT_JS)
        await browser.close()
        return total, dom


def build_image_pptx(out_dir: Path, total: int, pptx_out: Path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for i in range(1, total + 1):
        img = out_dir / f"slide-{i:02d}.png"
        if not img.exists():
            print(f"!! skip {img.name}")
            continue
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(
            str(img), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    prs.save(str(pptx_out))
    print(f"Saved (image)    : {pptx_out}")


def build_editable_pptx(dom, pptx_out: Path):
    """Reconstruct an editable deck from the extracted DOM."""
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_PX * PX_TO_EMU)
    prs.slide_height = Emu(SLIDE_H_PX * PX_TO_EMU)
    blank = prs.slide_layouts[6]

    def E(px):  # px -> EMU
        return Emu(int(round(px * PX_TO_EMU)))

    for sd in dom:
        s = prs.slides.add_slide(blank)
        # slide background
        bg = _rgb(sd.get("bg")) or DEFAULT_BG
        bgsp = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bgsp.fill.solid()
        bgsp.fill.fore_color.rgb = bg
        bgsp.line.fill.background()
        bgsp.shadow.inherit = False

        # boxes: largest first so small ones sit on top
        for b in sorted(sd.get("boxes", []), key=lambda b: -b["area"]):
            w, h = b["w"], b["h"]
            if w < 3 or h < 3:
                continue
            shape_type = (
                MSO_SHAPE.OVAL
                if b.get("oval")
                else (
                    MSO_SHAPE.ROUNDED_RECTANGLE
                    if (b.get("radius") or 0) > 1.5
                    else MSO_SHAPE.RECTANGLE
                )
            )
            sp = s.shapes.add_shape(shape_type, E(b["x"]), E(b["y"]), E(w), E(h))
            sp.shadow.inherit = False
            fill = _rgb(b.get("fill"))
            if fill is not None:
                sp.fill.solid()
                sp.fill.fore_color.rgb = fill
            else:
                sp.fill.background()
            border = _rgb(b.get("border"))
            if border is not None:
                sp.line.color.rgb = border
                sp.line.width = Pt(max(0.5, (b.get("borderW") or 1) * PX_TO_PT * 1.6))
            else:
                sp.line.fill.background()
            if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                try:
                    sp.adjustments[0] = max(
                        0.0, min(0.5, (b.get("radius") or 0) / min(w, h))
                    )
                except Exception:
                    pass

        # text boxes
        for t in sd.get("texts", []):
            tb = s.shapes.add_textbox(
                E(t["x"] - 2),
                E(t["y"] - 2),
                E(t["w"] + 6),
                E(max(t["h"] + 6, t["fontSize"] * 1.4)),
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            align = _align(t.get("align"))
            ratio = t.get("lineRatio") or 1.15
            # split runs on embedded newlines into paragraphs
            paras = [[]]
            for r in t.get("runs", []):
                parts = r["text"].split("\n")
                for k, part in enumerate(parts):
                    if k > 0:
                        paras.append([])
                    if part != "":
                        paras[-1].append(dict(r, text=part))
            first = True
            for prun in paras:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                p.alignment = align
                try:
                    p.line_spacing = ratio
                except Exception:
                    pass
                if not prun:
                    # empty line
                    rr = p.add_run()
                    rr.text = " "
                    rr.font.size = Pt(max(6, t["fontSize"] * PX_TO_PT * 0.6))
                    continue
                for r in prun:
                    rr = p.add_run()
                    rr.text = r["text"]
                    size_px = r.get("size") or t.get("fontSize") or 18
                    rr.font.size = Pt(max(6, round(size_px * PX_TO_PT, 1)))
                    rr.font.bold = bool(r.get("bold"))
                    fam = (
                        "Consolas"
                        if r.get("mono")
                        else (r.get("family") or "Noto Sans JP")
                    )
                    rr.font.name = fam
                    col = _rgb(r.get("color"))
                    if col is not None:
                        rr.font.color.rgb = col
                    _set_ea(rr, fam if not r.get("mono") else "Noto Sans JP")

    prs.save(str(pptx_out))
    print(f"Saved (editable) : {pptx_out}")


# --------------------------------------------------------------------------- #


def _editable_path(out: Path) -> Path:
    if out.stem.endswith("-editable"):
        return out
    return out.with_name(out.stem + "-editable" + out.suffix)


async def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("html", help="Path to slides HTML file")
    ap.add_argument(
        "output",
        help="Output PPTX file path (image deck). An '-editable' twin is also written.",
    )
    ap.add_argument(
        "--total",
        type=int,
        default=None,
        help="Number of slides (auto-detect if omitted)",
    )
    ap.add_argument(
        "--no-editable", action="store_true", help="Only produce the image-based PPTX"
    )
    ap.add_argument(
        "--only-editable",
        action="store_true",
        help="Only produce the text-editable PPTX",
    )
    args = ap.parse_args()

    html_path = Path(args.html)
    pptx_out = Path(args.output)
    out_dir = pptx_out.parent / "slides-png"
    out_dir.mkdir(exist_ok=True, parents=True)

    total, dom = await render_screenshots(html_path, out_dir, args.total)

    if not args.only_editable:
        build_image_pptx(out_dir, total, pptx_out)
    if not args.no_editable:
        build_editable_pptx(dom, _editable_path(pptx_out))


if __name__ == "__main__":
    asyncio.run(main())
